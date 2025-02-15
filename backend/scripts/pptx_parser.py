from pptx import Presentation
import sys
import os

def extract_text_from_pptx(pptx_path, output_path):
    try:
        # 加载PPT文件
        prs = Presentation(pptx_path)
        
        # 准备存储提取的文本
        slides_text = []
        
        # 遍历所有幻灯片
        for i, slide in enumerate(prs.slides, 1):
            slide_content = []
            slide_content.append(f"=== 幻灯片 {i} ===")
            
            # 提取幻灯片中的所有文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            
            slides_text.append("\n".join(slide_content))
        
        # 将所有文本写入输出文件
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n\n'.join(slides_text))
            
        print(f"成功提取文本到: {output_path}")
        return True
        
    except Exception as e:
        print(f"错误: {str(e)}")
        return False

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("使用方法: python pptx_parser.py <pptx文件路径> <输出文件路径>")
        sys.exit(1)
        
    pptx_path = sys.argv[1]
    output_path = sys.argv[2]
    
    success = extract_text_from_pptx(pptx_path, output_path)
    sys.exit(0 if success else 1)